# -*- coding: utf-8 -*-
"""核心文档解析服务。

依赖（硬）: trafilatura, lxml
依赖（软，缺失时部分格式降级）: pdfplumber, PyMuPDF (fitz), markitdown-no-magika,
    python-pptx, python-docx, openpyxl, xlrd, Pillow, beautifulsoup4
"""

import asyncio
import csv
import json
import logging
import os
import re
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO
from typing import Any, Dict, List, Optional

import httpx

from services.document_asset_service import DocumentAssetUploadService
from services.document_parse_models import DocumentAsset, DocumentParseResult, DocumentParseWarning

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    import fitz
except ImportError:
    fitz = None

import trafilatura
from trafilatura.settings import use_config

try:
    from markitdown import MarkItDown
except ImportError:
    MarkItDown = None
except Exception:
    MarkItDown = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE as PptxShapeType
except ImportError:
    PptxPresentation = None
    PptxShapeType = None

try:
    from openpyxl import load_workbook as openpyxl_load_workbook
except ImportError:
    openpyxl_load_workbook = None

try:
    import xlrd  # type: ignore[import-not-found]
except ImportError:
    xlrd = None

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

__all__ = [
    "EmbeddedImageUploader",
    "DataCleaningPipeline",
    "DocumentParserService",
]

# ==============================================================================
# ============== 嵌入图片上传服务 (EmbeddedImageUploader) ==============
# ==============================================================================

from utils.settings import settings as _settings  # noqa: E402  (settings 单点入口)

logger = logging.getLogger(__name__)


class EmbeddedImageUploader:
    """
    内嵌图片提取兼容类。

    上传职责已迁移到 DocumentAssetUploadService；本类保留 upload_images
    委托方法，是为了兼容旧调用方和现有图片提取逻辑。
    """
    MAX_BATCH = 10
    SUPPORTED_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif'}
    MIME_MAP = {
        '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
        '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp',
        '.tiff': 'image/tiff', '.tif': 'image/tiff', '.svg': 'image/svg+xml',
    }

    @classmethod
    def upload_images(cls, images: List[tuple]) -> Dict[str, str]:
        """
        批量上传图片，返回 {原始文件名: 可访问URL} 映射。
        images: [(filename, binary_data, mime_type), ...]
        """
        return DocumentAssetUploadService.upload_images(images)

    @classmethod
    def extract_from_zip(cls, data: bytes, media_prefix: str, min_size: int = 5120) -> List[tuple]:
        """从 ZIP 格式文档 (docx/pptx) 中提取 media 目录下的图片。"""
        images = []
        try:
            with zipfile.ZipFile(BytesIO(data)) as zf:
                for name in zf.namelist():
                    if not name.startswith(media_prefix):
                        continue
                    basename = os.path.basename(name)
                    ext_lower = os.path.splitext(basename)[1].lower()
                    if ext_lower not in cls.SUPPORTED_EXTS:
                        continue
                    img_data = zf.read(name)
                    if len(img_data) < min_size:
                        continue
                    mime = cls.MIME_MAP.get(ext_lower, 'image/png')
                    images.append((basename, img_data, mime))
        except Exception as e:
            logger.warning("zip image extraction failed: %s", e)
        return images

    @classmethod
    def extract_from_pdf(cls, data: bytes, max_pages: int = 50, min_size: int = 5120, min_dim: int = 50) -> List[tuple]:
        """从 PDF 中提取嵌入图片 (fitz)，按页面顺序返回。"""
        images = []
        if fitz is None:
            return images
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                img_idx = 0
                for pi in range(min(len(doc), max_pages)):
                    page = doc.load_page(pi)
                    page_dict = page.get_text("dict", sort=True)
                    for block in page_dict.get("blocks", []):
                        if block["type"] != 1:
                            continue
                        bbox = block.get("bbox", [0, 0, 0, 0])
                        w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
                        if w < min_dim or h < min_dim:
                            continue
                        img_bytes = block.get("image", b"")
                        if len(img_bytes) < min_size:
                            continue
                        img_idx += 1
                        ext = "png"
                        if img_bytes[:3] == b'\xff\xd8\xff':
                            ext = "jpg"
                        fname = f"pdf_image_{img_idx}.{ext}"
                        images.append((fname, img_bytes, f"image/{ext}"))
        except Exception as e:
            logger.warning("pdf image extraction failed: %s", e)
        return images
# ==============================================================================
# ================ 数据清洗管道 (DataCleaningPipeline) ================
# ==============================================================================

class DataCleaningPipeline:
    """LLM 友好的多阶段数据清洗管道，所有输出均为干净的 Markdown 字符串。"""

    _NOISY_PATTERNS = [re.compile(p, re.IGNORECASE) for p in [
        r'^[\-=*#_]{3,}$',
        r'.*\.(html|shtml|htm|php)\s*$',
        r'.{0,50}(搜狐|网易|腾讯|新浪|登录|注册|版权所有|版权声明).{0,50}$',
        r'\[\d+\]|\[下一页\]|\[上一页\]',
        r'\[(编辑|查看历史|讨论|阅读|来源|原标题)\]',
        r'^\*+\s*\[.*?\]\(.*?\)',
        r'^\s*(分享到|扫描二维码|返回搜狐|查看更多|责任编辑|记者|通讯员)',
        r'^\s*([京公网安备京网文京ICP备]|互联网新闻信息服务许可证|信息网络传播视听节目许可证)',
    ]]
    _IMG_PATTERN = re.compile(r'(!\[(.*?)\]\((.*?)\))')
    _LINK_PATTERN = re.compile(r'\[.*?\]\(.*?\)')
    _EDITOR_PATTERN = re.compile(r'(\(|\[)\s*责任编辑：.*?\s*(\)|\])')
    _PAGE_NUM_PATTERN = re.compile(
        r'^\s*[-—]\s*\d+\s*[-—]\s*$|'
        r'^\s*第\s*\d+\s*页\s*(共\s*\d+\s*页)?\s*$|'
        r'^\s*Page\s+\d+\s*(of\s+\d+)?\s*$',
        re.IGNORECASE
    )
    _REPEATED_LINE_THRESHOLD = 3

    def __init__(self, max_content_length: int = 80000):
        self.max_content_length = max_content_length

    @classmethod
    def _is_noisy_line(cls, line: str) -> bool:
        stripped = line.strip()
        if not stripped:
            return True
        for pat in cls._NOISY_PATTERNS:
            if pat.search(stripped):
                return True
        links = cls._LINK_PATTERN.findall(stripped)
        if len(links) > 2 and len(stripped) / (len(links) + 1) < 30:
            return True
        return False

    @staticmethod
    def _normalize_whitespace(text: str) -> str:
        lines = text.splitlines()
        out, prev_empty = [], False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not prev_empty:
                    out.append("")
                prev_empty = True
            else:
                out.append(stripped)
                prev_empty = False
        return "\n".join(out).strip()

    @classmethod
    def _remove_repeated_headers_footers(cls, text: str) -> str:
        lines = text.splitlines()
        if len(lines) < 20:
            return text
        line_counts: Dict[str, int] = {}
        for line in lines:
            s = line.strip()
            if s and len(s) < 100:
                line_counts[s] = line_counts.get(s, 0) + 1
        repeated = {s for s, c in line_counts.items() if c >= cls._REPEATED_LINE_THRESHOLD}
        if not repeated:
            return text
        return "\n".join(l for l in lines if l.strip() not in repeated)

    def _truncate(self, text: str, label: str = "内容") -> str:
        if len(text) > self.max_content_length:
            return text[:self.max_content_length] + f"\n\n...[{label}过长，已截断至 {self.max_content_length} 字符]"
        return text

    def clean_document(self, text: str) -> str:
        if not text:
            return ""
        text = self._remove_repeated_headers_footers(text)
        lines = text.splitlines()
        cleaned = []
        for line in lines:
            if self._PAGE_NUM_PATTERN.search(line.strip()):
                continue
            if self._is_noisy_line(line):
                continue
            line = self._EDITOR_PATTERN.sub('', line).strip()
            if line:
                cleaned.append(line)
        result = self._normalize_whitespace("\n".join(cleaned))
        return self._truncate(result, "文档内容")

    def clean_html(self, text: str) -> str:
        if not text:
            return ""
        lines = text.splitlines()
        cleaned = []
        for line in lines:
            if self._is_noisy_line(line):
                continue
            line = self._EDITOR_PATTERN.sub('', line).strip()
            if line:
                cleaned.append(line)
        result = self._normalize_whitespace("\n".join(cleaned))
        return self._truncate(result, "网页内容")

    def clean_table(self, text: str) -> str:
        if not text:
            return ""
        text = self._normalize_whitespace(text)
        return self._truncate(text, "表格内容")

    def clean_text(self, text: str) -> str:
        if not text:
            return ""
        text = self._normalize_whitespace(text)
        return self._truncate(text)

    async def validate_image_urls(self, md_text: str, client: httpx.AsyncClient) -> str:
        MAX_TO_VALIDATE = 25
        matches = list(self._IMG_PATTERN.finditer(md_text))
        if not matches:
            return md_text
        urls_all = {m.group(3).strip() for m in matches}
        urls_to_check = set(list(urls_all)[:MAX_TO_VALIDATE])

        async def _check(u):
            if not u or not u.startswith(('http://', 'https://')):
                return u, False
            try:
                resp = await client.head(u, timeout=5, follow_redirects=True)
                ct = resp.headers.get('content-type', '').lower()
                return u, resp.is_success and 'image' in ct
            except Exception:
                return u, False

        results = await asyncio.gather(*[_check(u) for u in urls_to_check], return_exceptions=True)
        valid = set()
        for r in results:
            if isinstance(r, tuple) and r[1]:
                valid.add(r[0])
        valid.update(urls_all - urls_to_check)

        def _replacer(m):
            return m.group(0) if m.group(3).strip() in valid else ""

        return self._IMG_PATTERN.sub(_replacer, md_text)
# ==============================================================================
# ============ 统一文件解析服务 (DocumentParserService) ============
# ==============================================================================

class DocumentParserService:
    """
    统一文件解析服务。
    支持: pdf, docx, doc, pptx, ppt, xlsx, xls, csv,
          html, htm, json, xml, txt, md,
          jpg, jpeg, png, gif, webp, bmp
    新入口返回 DocumentParseResult；parse/parse_async 保留字符串兼容行为。
    """

    def __init__(
            self,
            enable_embedded_image_upload: bool = True,
            enable_ocr: bool = True,
            pdf_max_pages: Optional[int] = None,
            max_content_chars: Optional[int] = None,
            max_table_rows: Optional[int] = None,
            min_img_bytes: Optional[int] = None,
            min_img_dim: Optional[int] = None,
    ):
        self.enable_embedded_image_upload = enable_embedded_image_upload
        self.enable_ocr = enable_ocr
        self.max_content_chars = max_content_chars or _settings.FILE_PARSE_MAX_CONTENT_CHARS
        self.cleaner = DataCleaningPipeline(max_content_length=self.max_content_chars)
        self.pdf_max_pages = pdf_max_pages or _settings.DOC_PARSER_PDF_MAX_PAGES
        self.max_table_rows = max_table_rows or _settings.DOC_PARSER_MAX_TABLE_ROWS
        self.min_img_bytes = min_img_bytes or _settings.DOC_PARSER_MIN_IMG_BYTES
        self.min_img_dim = min_img_dim or _settings.DOC_PARSER_MIN_IMG_DIM
        self._reset_parse_state("")
        self._markitdown = None
        if MarkItDown:
            try:
                self._markitdown = MarkItDown()
            except Exception as e:
                logger.warning("MarkItDown init failed: %s", e)

    def _reset_parse_state(self, ext: str) -> None:
        self._current_ext = ext
        self._meta: Dict[str, Any] = {}
        self._warnings: List[DocumentParseWarning] = []
        self._assets: List[DocumentAsset] = []
        self._parser_used = "DocumentParserService"
        self._fallback_used = False

    def _warn(self, code: str, message: str, level: str = "warning") -> None:
        if not any(w.code == code and w.message == message for w in self._warnings):
            self._warnings.append(DocumentParseWarning(code=code, message=message, level=level))

    def _mark_parser(self, parser_used: str, *, fallback: bool = False, reason: str = "") -> None:
        self._parser_used = parser_used
        self._fallback_used = self._fallback_used or fallback
        if fallback and reason:
            self._meta["fallback_reason"] = reason

    @staticmethod
    def _content_kind_from_ext(ext: str) -> str:
        if ext == ".pdf":
            return "pdf"
        if ext in (".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"):
            return "office"
        if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"):
            return "image"
        if ext in (".html", ".htm"):
            return "html"
        if ext in (".csv", ".txt", ".md", ".markdown", ".json", ".xml"):
            return "text"
        return "other"

    @staticmethod
    def _is_truncated_text(text: str) -> bool:
        return any(marker in text for marker in ("已截断", "内容过长", "JSON 内容过长", "XML 过长"))

    # ── MarkItDown 通用转换 ──────────────────────────────────
    def _markitdown_convert(self, data: bytes, suffix: str) -> str:
        if not self._markitdown:
            return ""
        if not suffix.startswith("."):
            suffix = f".{suffix}"
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
                f.write(data)
                tmp_path = f.name
            result = self._markitdown.convert(tmp_path)
            return result.text_content if result and result.text_content else ""
        except Exception as e:
            logger.warning("MarkItDown conversion failed for %s: %s", suffix, e)
            return ""
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass

    # ── PDF ──────────────────────────────────────────────────
    @staticmethod
    def _bbox_overlap(bbox_a, bbox_b, tolerance=2.0) -> bool:
        """判断两个 bbox 是否在 Y 轴方向上有足够重叠 (用于去重表格区域内的散碎文本)"""
        ax0, ay0, ax1, ay1 = bbox_a
        bx0, by0, bx1, by1 = bbox_b
        if ax1 < bx0 + tolerance or bx1 < ax0 + tolerance:
            return False
        if ay1 < by0 + tolerance or by1 < ay0 + tolerance:
            return False
        overlap_x = min(ax1, bx1) - max(ax0, bx0)
        width_a = ax1 - ax0
        return width_a > 0 and (overlap_x / width_a) > 0.5

    def _parse_pdf(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("pdfplumber+PyMuPDF")
        parts = []
        img_count = 0
        table_bboxes_per_page: Dict[int, list] = {}

        plumber_tables_per_page: Dict[int, list] = {}
        try:
            if pdfplumber is None:
                raise RuntimeError("pdfplumber not installed")
            with pdfplumber.open(BytesIO(data)) as plumber_pdf:
                page_limit = min(len(plumber_pdf.pages), self.pdf_max_pages)
                for pi in range(page_limit):
                    pp = plumber_pdf.pages[pi]
                    tables = pp.find_tables()
                    if not tables:
                        continue
                    page_tables = []
                    page_bboxes = []
                    for tbl in tables:
                        rows = tbl.extract()
                        if not rows:
                            continue
                        cleaned = []
                        for row in rows:
                            cleaned.append([(c or "").strip() for c in row])
                        if any(any(cell for cell in r) for r in cleaned):
                            page_tables.append((tbl.bbox[1], cleaned))
                            page_bboxes.append(tbl.bbox)
                    if page_tables:
                        plumber_tables_per_page[pi] = page_tables
                        table_bboxes_per_page[pi] = page_bboxes
        except Exception as e:
            self._warn("pdf_table_extraction_failed", "PDF 表格提取异常，正文解析会继续。")
            logger.warning("pdfplumber table extraction failed: %s", e)

        try:
            if fitz is None:
                raise RuntimeError("PyMuPDF not installed")
            with fitz.open(stream=data, filetype="pdf") as doc:
                total = len(doc)
                limit = min(total, self.pdf_max_pages)
                table_count = sum(len(v) for v in plumber_tables_per_page.values())
                self._meta.update({
                    "total_pages": total,
                    "processed_pages": limit,
                    "page_limit_applied": total > limit,
                    "table_count": table_count,
                })
                if total > self.pdf_max_pages:
                    self._warn(
                        "pdf_page_limit_applied",
                        f"PDF 共 {total} 页，已处理前 {limit} 页。",
                    )
                    logger.info("PDF has %s pages, only first %s pages processed", total, limit)

                for pi in range(limit):
                    page = doc.load_page(pi)
                    page_dict = page.get_text("dict", sort=True)
                    tbl_bboxes = table_bboxes_per_page.get(pi, [])
                    elements = []

                    for block in page_dict.get("blocks", []):
                        b_bbox = block.get("bbox", [0, 0, 0, 0])
                        y0 = b_bbox[1]

                        if block["type"] == 0:
                            if tbl_bboxes and any(self._bbox_overlap(b_bbox, tb) for tb in tbl_bboxes):
                                continue
                            lines_text = []
                            for ln in block.get("lines", []):
                                span_txt = "".join(s.get("text", "") for s in ln.get("spans", []))
                                if span_txt.strip():
                                    lines_text.append(span_txt.strip())
                            if lines_text:
                                elements.append((y0, "\n".join(lines_text)))

                        elif block["type"] == 1:
                            w, h = b_bbox[2] - b_bbox[0], b_bbox[3] - b_bbox[1]
                            if w < self.min_img_dim or h < self.min_img_dim:
                                continue
                            img_bytes = block.get("image", b"")
                            if len(img_bytes) < self.min_img_bytes:
                                continue
                            img_count += 1
                            elements.append(
                                (y0, f"![图片{img_count} (第{pi + 1}页, {int(w)}x{int(h)})](pdf_image_{img_count})"))

                    for tbl_y0, tbl_rows in plumber_tables_per_page.get(pi, []):
                        elements.append((tbl_y0, self._rows_to_md_table(tbl_rows)))

                    elements.sort(key=lambda x: x[0])
                    page_content = "\n\n".join(e[1] for e in elements)
                    if page_content.strip():
                        if limit > 1:
                            parts.append(f"<!-- 第 {pi + 1} 页 -->\n\n{page_content}")
                        else:
                            parts.append(page_content)

                if total > self.pdf_max_pages:
                    parts.append(f"\n\n> PDF 共 {total} 页，已处理前 {limit} 页")
                self._meta["embedded_image_count"] = img_count

            result = "\n\n".join(parts).strip()
            if result:
                result = self._upload_embedded_images(data, '.pdf', result)
                return self.cleaner.clean_document(result)
        except Exception as e:
            logger.warning("PDF PyMuPDF parse failed, fallback to MarkItDown: %s", e)
        fb = self._markitdown_convert(data, ".pdf")
        if not fb:
            return ""
        self._mark_parser("MarkItDown", fallback=True, reason="pdf_primary_parse_failed")
        self._warn("parser_fallback_used", "PDF 专用解析失败，已回退到 MarkItDown。")
        fb = self._upload_embedded_images(data, '.pdf', fb)
        return self.cleaner.clean_document(fb)

    # ── DOCX ─────────────────────────────────────────────────
    def _parse_docx(self, data: bytes, source_url: str = "") -> str:
        """DOCX: MarkItDown 优先 -> python-docx 回退(含图片) -> 清洗"""
        md_text = self._markitdown_convert(data, ".docx")
        if md_text:
            self._mark_parser("MarkItDown")
            md_text = self._upload_embedded_images(data, '.docx', md_text)
            return self.cleaner.clean_document(md_text)

        if DocxDocument is None:
            return ""
        try:
            self._mark_parser("python-docx", fallback=True, reason="markitdown_empty_result")
            self._warn("parser_fallback_used", "DOCX MarkItDown 解析无结果，已回退到 python-docx。")
            doc = DocxDocument(BytesIO(data))

            images = EmbeddedImageUploader.extract_from_zip(data, 'word/media/', min_size=self.min_img_bytes)
            url_map: Dict[str, str] = {}
            if images:
                images.sort(key=lambda x: x[0])
                logger.info("python-docx fallback extracted %s images", len(images))
                url_map = EmbeddedImageUploader.upload_images(images)
                if url_map:
                    logger.info("python-docx fallback uploaded %s/%s images", len(url_map), len(images))

            img_count = 0
            rId_to_url: Dict[str, str] = {}
            for rel_id, rel in doc.part.rels.items():
                if "image" in getattr(rel, 'reltype', ''):
                    target = os.path.basename(str(rel.target_ref))
                    if target.lower() in {k.lower() for k in url_map}:
                        for k, v in url_map.items():
                            if k.lower() == target.lower():
                                rId_to_url[rel_id] = v
                                break

            paragraphs = []
            for para in doc.paragraphs:
                para_xml = para._element.xml
                has_image = '<w:drawing' in para_xml or '<v:imagedata' in para_xml or '<wp:inline' in para_xml
                text = para.text.strip()
                if text:
                    if para.style and para.style.name and 'Heading' in para.style.name:
                        level = para.style.name.replace('Heading', '').strip()
                        prefix = '#' * (int(level) if level.isdigit() else 2)
                        paragraphs.append(f"{prefix} {text}")
                    else:
                        paragraphs.append(text)
                if has_image:
                    img_count += 1
                    img_url = None
                    embed_match = re.search(r'r:embed="([^"]+)"', para_xml)
                    if embed_match:
                        img_url = rId_to_url.get(embed_match.group(1))
                    if not img_url and url_map:
                        ordered = sorted(url_map.values())
                        idx = img_count - 1
                        if idx < len(ordered):
                            img_url = ordered[idx]
                    if img_url:
                        paragraphs.append(f"![文档图片{img_count}]({img_url})")

            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    paragraphs.append(self._rows_to_md_table(rows))

            md_text = "\n\n".join(paragraphs)
            self._meta.update({
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "embedded_image_count": len(images),
            })
            logger.info("python-docx fallback succeeded (%s chars)", len(md_text))
        except Exception as e:
            logger.warning("python-docx fallback failed: %s", e)
            return ""

        if not md_text:
            return ""
        return self.cleaner.clean_document(md_text)

    # ── PPTX ─────────────────────────────────────────────────
    def _parse_pptx(self, data: bytes, source_url: str = "") -> str:
        if PptxPresentation is not None:
            try:
                self._mark_parser("python-pptx")
                prs = PptxPresentation(BytesIO(data))
                parts = []
                img_count = 0
                table_count = 0
                images_to_upload: List[tuple] = []

                for si, slide in enumerate(prs.slides):
                    slide_title = ""
                    elements = []
                    for shape in slide.shapes:
                        top = shape.top or 0
                        if shape.has_text_frame:
                            text = "\n".join(
                                p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                            )
                            if text:
                                if not slide_title:
                                    try:
                                        if shape.is_placeholder and shape.placeholder_format.idx == 0:
                                            slide_title = text
                                    except Exception:
                                        pass
                                elements.append((top, text))
                        if PptxShapeType and shape.shape_type == PptxShapeType.PICTURE:
                            img_count += 1
                            placeholder = f"__PPTX_IMG_{img_count}__"
                            try:
                                blob = shape.image.blob
                                ext = getattr(shape.image, 'ext', 'png') or 'png'
                                if len(blob) >= self.min_img_bytes:
                                    fname = f"pptx_s{si + 1}_img{img_count}.{ext}"
                                    mime = f"image/{ext}"
                                    images_to_upload.append((placeholder, fname, blob, mime))
                            except Exception:
                                pass
                            elements.append((top, f"![幻灯片{si + 1}-图片{img_count}]({placeholder})"))
                        if shape.has_table:
                            rows_data = []
                            for row in shape.table.rows:
                                rows_data.append([cell.text.strip() for cell in row.cells])
                            if rows_data:
                                table_count += 1
                                elements.append((top, self._rows_to_md_table(rows_data)))
                    elements.sort(key=lambda x: x[0])
                    header = f"## 幻灯片 {si + 1}"
                    if slide_title:
                        header += f": {slide_title}"
                    body = "\n\n".join(e[1] for e in elements)
                    if body.strip():
                        parts.append(f"{header}\n\n{body}")

                if parts:
                    md_text = "\n\n---\n\n".join(parts)
                    if images_to_upload:
                        upload_list = [(fn, bl, mi) for _, fn, bl, mi in images_to_upload]
                        logger.info("PPTX extracted %s images", len(upload_list))
                        url_map = EmbeddedImageUploader.upload_images(upload_list)
                        if url_map:
                            logger.info("PPTX uploaded %s/%s images", len(url_map), len(upload_list))
                            for ph, fn, _, _ in images_to_upload:
                                if fn in url_map:
                                    md_text = md_text.replace(f"]({ph})", f"]({url_map[fn]})")
                    self._meta.update({
                        "slide_count": len(prs.slides),
                        "table_count": table_count,
                        "embedded_image_count": len(images_to_upload),
                    })
                    return self.cleaner.clean_document(md_text)
            except Exception as e:
                logger.warning("python-pptx parse failed, fallback to MarkItDown: %s", e)
        fb = self._markitdown_convert(data, ".pptx")
        if not fb:
            return ""
        self._mark_parser("MarkItDown", fallback=True, reason="python_pptx_failed")
        self._warn("parser_fallback_used", "PPTX 专用解析失败，已回退到 MarkItDown。")
        fb = self._upload_embedded_images(data, '.pptx', fb)
        return self.cleaner.clean_document(fb)

    # ── Excel (xlsx / xls) ───────────────────────────────────
    def _parse_excel(self, data: bytes, source_url: str = "") -> str:
        is_xls = not data[:4] == b'PK\x03\x04'
        suffix = ".xls" if is_xls else ".xlsx"
        md_text = self._markitdown_convert(data, suffix)
        if md_text and md_text.strip():
            self._mark_parser("MarkItDown")
            return self.cleaner.clean_table(md_text)
        if not is_xls and openpyxl_load_workbook is not None:
            try:
                self._mark_parser("openpyxl", fallback=True, reason="markitdown_empty_result")
                wb = openpyxl_load_workbook(BytesIO(data), read_only=True, data_only=True)
                parts = []
                row_limit_applied = False
                processed_rows = 0
                for name in wb.sheetnames:
                    ws = wb[name]
                    rows = []
                    for ri, row in enumerate(ws.iter_rows(values_only=True)):
                        if ri >= self.max_table_rows:
                            rows.append(["...", f"共 {ws.max_row} 行，已截断", "..."])
                            row_limit_applied = True
                            break
                        rows.append([str(c) if c is not None else "" for c in row])
                        processed_rows += 1
                    if rows:
                        parts.append(f"### 工作表: {name}\n\n{self._rows_to_md_table(rows)}")
                self._meta.update({
                    "sheet_count": len(wb.sheetnames),
                    "processed_rows": processed_rows,
                    "row_limit_applied": row_limit_applied,
                })
                if row_limit_applied:
                    self._warn("table_row_limit_applied", f"表格最多保留前 {self.max_table_rows} 行。")
                wb.close()
                if parts:
                    return self.cleaner.clean_table("\n\n".join(parts))
            except Exception as e:
                logger.warning("openpyxl parse failed: %s", e)
        if is_xls and xlrd is not None:
            try:
                self._mark_parser("xlrd", fallback=True, reason="markitdown_empty_result")
                wb = xlrd.open_workbook(file_contents=data)
                parts = []
                row_limit_applied = False
                processed_rows = 0
                for name in wb.sheet_names():
                    ws = wb.sheet_by_name(name)
                    rows = []
                    for ri in range(min(ws.nrows, self.max_table_rows)):
                        rows.append([str(ws.cell_value(ri, ci)) for ci in range(ws.ncols)])
                        processed_rows += 1
                    if ws.nrows > self.max_table_rows:
                        rows.append(["...", f"共 {ws.nrows} 行，已截断", "..."])
                        row_limit_applied = True
                    if rows:
                        parts.append(f"### 工作表: {name}\n\n{self._rows_to_md_table(rows)}")
                self._meta.update({
                    "sheet_count": len(wb.sheet_names()),
                    "processed_rows": processed_rows,
                    "row_limit_applied": row_limit_applied,
                })
                if row_limit_applied:
                    self._warn("table_row_limit_applied", f"表格最多保留前 {self.max_table_rows} 行。")
                if parts:
                    return self.cleaner.clean_table("\n\n".join(parts))
            except Exception as e:
                logger.warning("xlrd parse failed: %s", e)
        return ""

    # ── CSV ──────────────────────────────────────────────────
    def _parse_csv(self, data: bytes, source_url: str = "") -> str:
        md_text = self._markitdown_convert(data, ".csv")
        if md_text and md_text.strip():
            self._mark_parser("MarkItDown")
            return self.cleaner.clean_table(md_text)
        text = self._decode_bytes(data)
        if not text:
            return ""
        self._mark_parser("csv")
        try:
            dialect = csv.Sniffer().sniff(text[:8192])
            reader = csv.reader(text.splitlines(), dialect)
        except csv.Error:
            reader = csv.reader(text.splitlines())
        rows = []
        row_limit_applied = False
        for i, row in enumerate(reader):
            if i >= self.max_table_rows:
                rows.append(["...", "[已截断]", "..."])
                row_limit_applied = True
                break
            rows.append(row)
        self._meta.update({
            "processed_rows": min(len(rows), self.max_table_rows),
            "row_limit_applied": row_limit_applied,
        })
        if row_limit_applied:
            self._warn("table_row_limit_applied", f"CSV 最多保留前 {self.max_table_rows} 行。")
        return self.cleaner.clean_table(self._rows_to_md_table(rows)) if rows else ""

    # ── HTML (文件) ──────────────────────────────────────────
    def _parse_html_file(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("trafilatura")
        text = self._decode_bytes(data)
        if not text:
            return ""
        cfg = use_config()
        cfg.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")
        result = trafilatura.extract(
            text, config=cfg, output_format='markdown',
            include_images=True, favor_recall=True
        )
        return self.cleaner.clean_html(result) if result else ""

    # ── JSON ─────────────────────────────────────────────────
    def _parse_json(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("json")
        text = self._decode_bytes(data)
        if not text:
            return ""
        try:
            obj = json.loads(text)
            formatted = json.dumps(obj, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            formatted = text
        if len(formatted) > self.max_content_chars:
            formatted = formatted[:self.max_content_chars] + "\n... [JSON 内容过长，已截断]"
            self._warn("json_content_truncated", f"JSON 内容过长，已截断至 {self.max_content_chars} 字符。")
            self._meta["content_limit_applied"] = True
        return f"```json\n{formatted}\n```"

    # ── XML ──────────────────────────────────────────────────
    def _parse_xml(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("xml")
        md_text = self._markitdown_convert(data, ".xml")
        if md_text and md_text.strip():
            self._mark_parser("MarkItDown")
            return self.cleaner.clean_text(md_text)
        text = self._decode_bytes(data)
        if not text:
            return ""
        if len(text) > self.max_content_chars:
            text = text[:self.max_content_chars] + "\n... [XML 过长，已截断]"
            self._warn("xml_content_truncated", f"XML 内容过长，已截断至 {self.max_content_chars} 字符。")
            self._meta["content_limit_applied"] = True
        return f"```xml\n{text}\n```"

    # ── Plain Text ───────────────────────────────────────────
    def _parse_plain_text(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("plain-text")
        text = self._decode_bytes(data)
        return self.cleaner.clean_text(text) if text else ""

    # ── Markdown ─────────────────────────────────────────────
    def _parse_markdown(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("markdown")
        text = self._decode_bytes(data)
        return self.cleaner.clean_text(text) if text else ""

    # ── Image ────────────────────────────────────────────────
    def _parse_image(self, data: bytes, source_url: str = "") -> str:
        self._mark_parser("Pillow")
        parts = []
        if source_url:
            parts.append(f"![image]({source_url})")
        if PILImage is not None:
            try:
                img = PILImage.open(BytesIO(data))
                w, h = img.size
                fmt = img.format or "Unknown"
                self._meta.update({
                    "width": w,
                    "height": h,
                    "format": fmt,
                    "mode": img.mode,
                })
                parts.append(f"**图片信息**: {fmt}, {w}x{h}px, {img.mode}")
                if self.enable_ocr:
                    try:
                        import pytesseract  # type: ignore[import-not-found]
                        ocr_text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                        if ocr_text and ocr_text.strip():
                            self._meta["ocr_applied"] = True
                            parts.append(f"\n**OCR 识别文本**:\n\n{ocr_text.strip()}")
                    except (ImportError, Exception):
                        self._warn("ocr_unavailable", "图片 OCR 未执行或执行失败。", level="info")
            except Exception as e:
                logger.warning("image parse failed: %s", e)
        if not parts:
            return f"[图片文件, {len(data)} bytes]"
        return "\n\n".join(parts)

    # ── 工具方法 ─────────────────────────────────────────────
    @staticmethod
    def _decode_bytes(data: bytes) -> str:
        for enc in ('utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1'):
            try:
                return data.decode(enc)
            except (UnicodeDecodeError, LookupError):
                continue
        return data.decode('utf-8', errors='replace')

    @staticmethod
    def _rows_to_md_table(rows: list) -> str:
        if not rows:
            return ""
        max_cols = max(len(r) for r in rows)
        padded = [r + [""] * (max_cols - len(r)) for r in rows]
        header = "| " + " | ".join(str(c).replace("|", "\\|").replace("\n", " ")[:80] for c in padded[0]) + " |"
        sep = "| " + " | ".join("---" for _ in padded[0]) + " |"
        body = []
        for row in padded[1:]:
            body.append("| " + " | ".join(str(c).replace("|", "\\|").replace("\n", " ")[:80] for c in row) + " |")
        return "\n".join([header, sep] + body)

    # ── 嵌入图片: 提取 + 上传 + 替换 ─────────────────────────
    _IMG_LOCAL_REF = re.compile(r'(!\[[^\]]*\])\((?!https?://|data:)([^)]+)\)')

    def _strip_and_replace_data_uris(self, md_text: str, ordered_urls: List[tuple]) -> str:
        """
        扫描 markdown, 把所有 ![alt](data:image/...;base64,...) 替换为上传后的真实 URL。
        完全不解码 base64 — 只用字符串定位 data: 开头和 ) 结尾, 然后整段替换。
        ordered_urls: [(filename, url), ...] 按文档中图片出现顺序排列。
        """
        url_idx = 0
        parts: List[str] = []
        pos = 0
        while pos < len(md_text):
            img_start = md_text.find('![', pos)
            if img_start == -1:
                parts.append(md_text[pos:])
                break
            bracket_close = md_text.find('](', img_start + 2)
            if bracket_close == -1:
                parts.append(md_text[pos:])
                break
            uri_start = bracket_close + 2
            if not md_text[uri_start:uri_start + 5] == 'data:':
                parts.append(md_text[pos:uri_start])
                pos = uri_start
                continue
            paren_close = md_text.find(')', uri_start)
            if paren_close == -1:
                parts.append(md_text[pos:])
                break
            alt = md_text[img_start + 2:bracket_close]
            parts.append(md_text[pos:img_start])
            if url_idx < len(ordered_urls):
                fname, url = ordered_urls[url_idx]
                parts.append(f"![{alt or fname}]({url})")
                url_idx += 1
            else:
                parts.append(f"![{alt or '图片'}]")
            pos = paren_close + 1
        return ''.join(parts)

    def _upload_embedded_images(self, data: bytes, ext: str, md_text: str) -> str:
        """
        核心逻辑: 从文档二进制直接提取图片 -> 上传到服务器获取真实 URL ->
        替换 markdown 中所有 data:image base64 引用和本地文件名引用。
        绝不解码 base64, 图片来源是文档 ZIP/PDF 二进制本身。
        """
        if not md_text:
            return md_text

        if not self.enable_embedded_image_upload:
            return self._strip_and_replace_data_uris(md_text, [])

        # ── Step 1: 从文档二进制提取真实图片文件 ──────────────────
        images: List[tuple] = []
        if ext == '.pdf':
            images = EmbeddedImageUploader.extract_from_pdf(
                data, max_pages=self.pdf_max_pages,
                min_size=self.min_img_bytes, min_dim=self.min_img_dim
            )
        elif ext in ('.docx', '.doc'):
            images = EmbeddedImageUploader.extract_from_zip(data, 'word/media/', min_size=self.min_img_bytes)
        elif ext in ('.pptx', '.ppt'):
            images = EmbeddedImageUploader.extract_from_zip(data, 'ppt/media/', min_size=self.min_img_bytes)

        if not images:
            return self._strip_and_replace_data_uris(md_text, [])

        images.sort(key=lambda x: x[0])
        asset_by_name: Dict[str, DocumentAsset] = {}
        for fname, img_data, mime in images:
            asset = DocumentAsset(
                filename=fname,
                mime_type=mime,
                size=len(img_data),
                status="extracted",
            )
            self._assets.append(asset)
            asset_by_name[fname] = asset

        self._meta["embedded_image_count"] = max(
            int(self._meta.get("embedded_image_count", 0) or 0),
            len(images),
        )

        if not _settings.DOC_PARSER_IMAGE_UPLOAD_URL:
            for asset in asset_by_name.values():
                asset.status = "skipped"
                asset.error = "image_upload_not_configured"
            self._warn("image_upload_not_configured", "未配置 DOC_PARSER_IMAGE_UPLOAD_URL，内嵌图片未上传。")
            return self._strip_and_replace_data_uris(md_text, [])

        # ── Step 2: 批量上传到服务器 ─────────────────────────────
        logger.info("extracted %s embedded images, uploading", len(images))
        url_map = EmbeddedImageUploader.upload_images(images)
        if not url_map:
            for asset in asset_by_name.values():
                asset.status = "failed"
                asset.error = "image_upload_failed"
            self._warn("image_upload_failed", "内嵌图片上传失败，已移除 base64 图片噪音。")
            logger.warning("embedded image upload failed, base64 noise will be stripped")
            return self._strip_and_replace_data_uris(md_text, [])

        ordered_urls = [(fname, url_map[fname]) for fname, _, _ in images if fname in url_map]
        for fname, url in ordered_urls:
            if fname in asset_by_name:
                asset_by_name[fname].status = "uploaded"
                asset_by_name[fname].url = url
        logger.info("uploaded %s/%s embedded images", len(ordered_urls), len(images))

        # ── Step 3: 替换 data:image base64 引用 (按顺序匹配) ────
        md_text = self._strip_and_replace_data_uris(md_text, ordered_urls)

        # ── Step 4: 替换本地文件名引用 ![](image1.png) (按名称匹配)
        local_matches = list(self._IMG_LOCAL_REF.finditer(md_text))
        if not local_matches:
            return md_text

        url_map_lower = {k.lower(): v for k, v in url_map.items()}
        name_no_ext_map = {os.path.splitext(k)[0].lower(): v for k, v in url_map.items()}

        for m_obj in reversed(local_matches):
            prefix = m_obj.group(1)
            ref = m_obj.group(2)
            ref_base = os.path.basename(ref).lower()
            ref_no_ext = os.path.splitext(ref_base)[0]
            new_url = (
                    url_map_lower.get(ref_base)
                    or url_map_lower.get(ref.lower())
                    or name_no_ext_map.get(ref_no_ext)
                    or name_no_ext_map.get(ref.lower())
            )
            if new_url:
                replacement = f"{prefix}({new_url})"
                md_text = md_text[:m_obj.start()] + replacement + md_text[m_obj.end():]

        return md_text

    # ── 主入口 (同步, 在 asyncio.to_thread 中调用) ───────────
    def _parse_markdown_core(self, binary_content: bytes, file_extension: str, source_url: str = "") -> str:
        ext = file_extension.lower().strip()
        if not ext.startswith("."):
            ext = f".{ext}"

        result = ""
        if ext == '.pdf':
            result = self._parse_pdf(binary_content, source_url)
        elif ext == '.docx':
            result = self._parse_docx(binary_content, source_url)
        elif ext == '.pptx':
            result = self._parse_pptx(binary_content, source_url)
        elif ext in ('.xlsx', '.xls'):
            result = self._parse_excel(binary_content, source_url)
        elif ext == '.csv':
            result = self._parse_csv(binary_content, source_url)
        elif ext in ('.html', '.htm'):
            result = self._parse_html_file(binary_content, source_url)
        elif ext == '.json':
            result = self._parse_json(binary_content, source_url)
        elif ext == '.xml':
            result = self._parse_xml(binary_content, source_url)
        elif ext == '.txt':
            result = self._parse_plain_text(binary_content, source_url)
        elif ext in ('.md', '.markdown'):
            result = self._parse_markdown(binary_content, source_url)
        elif ext in ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'):
            result = self._parse_image(binary_content, source_url)
        elif ext in ('.doc', '.ppt'):
            self._mark_parser("MarkItDown")
            fb = self._markitdown_convert(binary_content, ext)
            if fb:
                fb = self._upload_embedded_images(binary_content, ext, fb)
                result = self.cleaner.clean_document(fb)

        if not (result and result.strip()):
            logger.warning("specific parser returned empty for %s, trying MarkItDown fallback", ext)
            fb = self._markitdown_convert(binary_content, ext)
            if fb:
                if ext in ('.pdf', '.docx', '.doc', '.pptx', '.ppt'):
                    fb = self._upload_embedded_images(binary_content, ext, fb)
                self._mark_parser("MarkItDown", fallback=True, reason=f"{ext}_specific_parse_empty")
                self._warn("parser_fallback_used", f"格式 {ext} 专用解析无结果，已尝试 MarkItDown。")
                result = self.cleaner.clean_document(fb)
            else:
                result = f"[无法解析 {ext} 格式文件]"

        return result

    def parse_document(self, binary_content: bytes, file_extension: str, source_url: str = "") -> DocumentParseResult:
        ext = file_extension.lower().strip()
        if not ext.startswith("."):
            ext = f".{ext}"
        self._reset_parse_state(ext)
        markdown = self._parse_markdown_core(binary_content, ext, source_url)
        truncated = self._is_truncated_text(markdown)
        if truncated:
            self._warn("content_truncated", "解析结果已在核心解析层截断。")
            self._meta["content_truncated"] = True
        return DocumentParseResult(
            markdown=markdown,
            content_kind=self._content_kind_from_ext(ext),
            parser_used=self._parser_used,
            fallback_used=self._fallback_used,
            meta=dict(self._meta),
            warnings=list(self._warnings),
            assets=list(self._assets),
            truncated=truncated,
        )

    def parse(self, binary_content: bytes, file_extension: str, source_url: str = "") -> str:
        return self.parse_document(binary_content, file_extension, source_url).markdown

    async def parse_document_async(
            self,
            binary_content: bytes,
            file_extension: str,
            source_url: str = "",
    ) -> DocumentParseResult:
        return await asyncio.to_thread(self.parse_document, binary_content, file_extension, source_url)

    async def parse_async(self, binary_content: bytes, file_extension: str, source_url: str = "") -> str:
        return await asyncio.to_thread(self.parse, binary_content, file_extension, source_url)

    def parse_html_content(self, html: str, base_url: str = "") -> str:
        if not html:
            return ""
        cfg = use_config()
        cfg.set("DEFAULT", "EXTRACTION_TIMEOUT", "10")
        result = trafilatura.extract(
            html, config=cfg, output_format='markdown',
            include_images=True, favor_recall=True
        )
        return self.cleaner.clean_html(result) if result else ""
